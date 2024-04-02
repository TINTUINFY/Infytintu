from pptx import Presentation
import glob

class PptWrapper():
    """Can extract text from .pptx, .pptm, .pptb formats only"""
    
    def __init__(self, file):
        self.prs = Presentation(file)

    def read_ppt(self):  
        text =[]
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text+". ")

        result = "".join(text)
        result = result.replace("\n", "")
        return result